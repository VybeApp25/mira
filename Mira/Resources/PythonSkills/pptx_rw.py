#!/usr/bin/env python3
"""
Create and read PowerPoint (.pptx) presentations using python-pptx.

Args JSON (stdin):
  {op: "create",
   path: "<abs_path>.pptx",
   title?: "<presentation title>",
   slides: [
     {title: "<slide title>", body?: "<text or bullet list>",
      layout?: "title"|"title_content"|"blank",
      bullets?: ["point 1", "point 2", ...]},
     ...
   ]}
      → create a presentation; returns {success, path, slide_count}

  {op: "read", path: "<abs_path>.pptx"}
      → extract text from each slide; returns {success, slides:[{index, title, text}], total}

  {op: "append_slide",
   path: "<abs_path>.pptx",
   title: "<slide title>",
   body?: "<text>",
   bullets?: ["...", ...]}
      → add a slide to an existing presentation

Output JSON: {success, ...fields, error?}
"""
import sys, json, os

DEPS = ["python-pptx"]

def out(data: dict):
    print(json.dumps(data))

def load_pptx():
    try:
        import pptx
        return pptx
    except ImportError:
        out({"success": False, "error": "python-pptx not installed. PythonSkillRunner should install it automatically."})
        sys.exit(0)

# ── helpers ───────────────────────────────────────────────────────────────────

def expand(path: str) -> str:
    return os.path.expanduser(path)

def ensure_pptx_ext(path: str) -> str:
    return path if path.lower().endswith(".pptx") else path + ".pptx"

def add_slide(prs, slide_data: dict):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    layout_name = slide_data.get("layout", "title_content")
    slide_title = slide_data.get("title", "")
    body_text   = slide_data.get("body", "")
    bullets     = slide_data.get("bullets", [])

    # Layout indices: 0=title slide, 1=title+content, 5=title only, 6=blank
    layout_map = {"title": 0, "title_content": 1, "title_only": 5, "blank": 6}
    layout_idx = layout_map.get(layout_name, 1)
    # Clamp to valid range
    layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)

    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # Set title
    if slide_title and slide.shapes.title:
        slide.shapes.title.text = slide_title

    # Set body / bullets
    content = bullets if bullets else ([body_text] if body_text else [])
    if content:
        # Find content placeholder (index 1) or any text frame that isn't the title
        body_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                body_ph = ph
                break
        if body_ph and hasattr(body_ph, "text_frame"):
            tf = body_ph.text_frame
            tf.word_wrap = True
            for i, line in enumerate(content):
                if i == 0:
                    tf.text = line
                else:
                    tf.add_paragraph().text = line

# ── create ────────────────────────────────────────────────────────────────────

def create_presentation(path: str, title: str | None, slides: list):
    pptx = load_pptx()
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width  = pptx.util.Emu(9144000)
    prs.slide_height = pptx.util.Emu(5143500)

    # If a presentation title is given, insert a title slide first
    if title:
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = ""

    for slide_data in slides:
        add_slide(prs, slide_data)

    path = ensure_pptx_ext(expand(path))
    os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
    prs.save(path)
    out({"success": True, "path": path, "slide_count": len(prs.slides)})

# ── read ─────────────────────────────────────────────────────────────────────

def read_presentation(path: str):
    load_pptx()
    from pptx import Presentation

    path = expand(path)
    if not os.path.exists(path):
        out({"success": False, "error": f"File not found: {path}"})
        return

    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                title = text
            else:
                texts.append(text)
        slides.append({"index": i, "title": title, "text": "\n".join(texts)})

    out({"success": True, "slides": slides, "total": len(slides),
         "path": path})

# ── append_slide ─────────────────────────────────────────────────────────────

def append_slide_to_file(path: str, slide_data: dict):
    load_pptx()
    from pptx import Presentation

    path = ensure_pptx_ext(expand(path))
    if not os.path.exists(path):
        out({"success": False, "error": f"File not found: {path}. Use op=create first."})
        return

    prs = Presentation(path)
    add_slide(prs, slide_data)
    prs.save(path)
    out({"success": True, "path": path, "slide_count": len(prs.slides)})

# ── main ─────────────────────────────────────────────────────────────────────

def main():
    try:
        args = json.load(sys.stdin)
    except Exception:
        args = {}
    op = args.get("op", "create")

    if op == "create":
        path = args.get("path", "")
        if not path:
            out({"success": False, "error": "path is required"}); return
        slides = args.get("slides", [])
        create_presentation(path, args.get("title"), slides)

    elif op == "read":
        path = args.get("path", "")
        if not path:
            out({"success": False, "error": "path is required"}); return
        read_presentation(path)

    elif op == "append_slide":
        path = args.get("path", "")
        if not path:
            out({"success": False, "error": "path is required"}); return
        append_slide_to_file(path, {
            "title":   args.get("title", ""),
            "body":    args.get("body", ""),
            "bullets": args.get("bullets", []),
            "layout":  args.get("layout", "title_content"),
        })

    else:
        out({"success": False,
             "error": f"Unknown op: {op}. Use: create, read, append_slide"})

main()
